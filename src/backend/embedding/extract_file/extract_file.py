import os
import logging
import magic  # Used to detect file type
import pytesseract
from src.backend.minio.minio_manager import MinioManager
import pdfplumber
import json
import io
from docx import Document
from PIL import Image
from enum import Enum
import zipfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import openpyxl
from openpyxl.drawing.image import Image as OpenpyxlImage
from openpyxl.chart import (
    BarChart, PieChart, LineChart, AreaChart, RadarChart, ScatterChart, BubbleChart, DoughnutChart, StockChart,
    SurfaceChart
)
from bs4 import BeautifulSoup
import requests
from bs4.element import Comment
from io import BytesIO
import re
import yaml
from src.backend.embedding.lib.download_and_upload_file import LocalFileDownloadAndUpload

from src.backend.lib.singleton_class import Singleton

from src.backend.lib.logging_config import get_primitivechat_logger


# Configure Logging
logger = get_primitivechat_logger(__name__)

# Enum for file types
class FileType(Enum):
    PDF = ".pdf"
    DOCX = ".docx"
    PPTX = ".pptx"
    XLSX = ".xlsx"
    HTML = ".html"
    JSON = ".json"
    IMAGE = (".jpg", ".jpeg", ".png")
    YAML = ".yaml"
    CODE = (".java", ".py", ".cpp", ".c", ".php", ".js")

class CustomShapeType(Enum):
    PICTURE = 13

class UploadFileForChunks(metaclass=Singleton):
    def __init__(self):
        self.minio_manager = MinioManager()
        self.file_extract = FileExtractor()
        self.download_and_upload = LocalFileDownloadAndUpload()

    def extract_file(self, customer_guid: str, filename: str,local_path, test_mode: bool = None):
        logger.info(f"Extracting file '{filename}' for customer '{customer_guid}'")
        use_test_mode = test_mode

        # Skip download in test mode if local_path is provided
        if not use_test_mode:
            local_path = self.download_and_upload.download_and_save_file(customer_guid, filename)
        else:
            logger.info(f"Test mode: Skipping Minio download. Using local path: {local_path}")

        #Verify file type
        try:
            file_type = self.file_extract.detect_file_type(local_path)
            if file_type == FileType.PDF:
                logger.info(f"Verified file '{filename}' as a valid PDF.")
                output_file_path = self.file_extract.extract_pdf_content(customer_guid, local_path, filename)
            elif file_type == FileType.DOCX:
                logger.info(f"Verified file '{filename}' as a valid DOCX")
                output_file_path = self.file_extract.extract_docx_content(customer_guid, local_path, filename)
            elif file_type==FileType.PPTX:
                logger.info(f"Verified file '{filename}' as a valid PPTX.")
                output_file_path = self.file_extract.extract_ppt_content(customer_guid, local_path, filename)
            elif file_type == FileType.XLSX:
                logger.info(f"Verified file '{filename}' as a valid XLSX.")
                output_file_path = self.file_extract.extract_excel_content(customer_guid, local_path, filename)
            elif file_type == FileType.HTML:
                logger.info(f"Verified file '{filename}' as a valid HTML")
                output_file_path = self.file_extract.extract_html_content(customer_guid, local_path, filename)
            elif file_type == FileType.JSON:
                logger.info(f"Verified file '{filename}' as a valid JSON")
                output_file_path = self.file_extract.extract_json_content(customer_guid, local_path, filename)
            elif file_type == FileType.IMAGE:
                logger.info(f"Verified file '{filename}' as a valid JPEG or JPG or PNG")
                output_file_path = self.file_extract.extract_image_content(customer_guid, local_path, filename)
            elif file_type == FileType.YAML:
                logger.info(f"Verified file '{filename}' as a valid YAML or YML")
                output_file_path = self.file_extract.extract_yaml_content(customer_guid, local_path, filename)
            elif file_type == FileType.CODE:
                logger.info(f"Verified file '{filename}' as a valid CODE file.")
                output_file_path = self.file_extract.extract_code_file_content(customer_guid, local_path, filename)
            else:
                logger.error(f"File '{filename}' is not a valid file (detected type: {file_type})")
                raise Exception("Invalid file type: Uploaded file is not a Valid.")
        except Exception as e:
            logger.error(f"Error extracting file for '{filename}': {e}")
            raise Exception("File extraction failed.")

        #Upload extracted content
        try:
            if use_test_mode:
                return output_file_path
            else:    
                self.download_and_upload.upload_extracted_content(customer_guid, filename, output_file_path)
                return {"message": f"{file_type.name} extracted and uploaded successfully."}
        except Exception as e:
            logger.error(f"File upload error: {e}")
            raise Exception(f"File upload failed.{e}")

    def extract_html_files(self, customer_guid: str, source_url_link: str, filename: str):

        try:

            logger.info(f"Fetching HTML content from URL: {source_url_link}")
            response = requests.get(source_url_link)
            response.raise_for_status()
            html_content = response.text

            local_path = f"/tmp/{customer_guid}/{filename}"
            os.makedirs(os.path.dirname(local_path), exist_ok=True)
            with open(local_path, "w", encoding="utf-8") as html_file:
                html_file.write(html_content)

            logger.info(f"HTML content from '{source_url_link}' saved to '{local_path}'")

            output_file_path=self.file_extract.extract_html_content(customer_guid, local_path, filename)
            logger.info(f"HTML content from URL '{source_url_link}' extracted successfully.")

            self.download_and_upload.upload_extracted_content(customer_guid, filename, output_file_path)
            return {"message": "HTML extracted and uploaded successfully."}

        except requests.exceptions.RequestException as req_err:
            logger.error(f"Failed to fetch HTML content from URL '{source_url_link}': {req_err}")
            raise Exception(f"Failed to fetch HTML content: {req_err}")
        except Exception as e:
            logger.error(f"Error processing HTML file from URL '{source_url_link}': {e}")
            raise Exception(f"HTML file processing failed: {e}")

class FileExtractor:

    def detect_file_type(self, file_path: str):
        MIME_TO_EXTENSION = {
            "application/pdf": FileType.PDF,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": FileType.DOCX,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": FileType.PPTX,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": FileType.XLSX,
            "text/html": FileType.HTML,
            "application/xhtml+xml": FileType.HTML,
            "application/json": FileType.JSON,
            "text/json": FileType.JSON,
            "image/jpeg": FileType.IMAGE,
            "image/png": FileType.IMAGE,
            "image/jpg": FileType.IMAGE,
            "text/x-yaml": FileType.YAML,
            "application/x-yaml": FileType.YAML,
            "text/plain": FileType.YAML,
            "text/x-java-source": FileType.CODE,
            "text/x-java": FileType.CODE,
            "text/x-python": FileType.CODE,
            "text/x-script.python": FileType.CODE,
            "text/x-c": FileType.CODE,
            "text/x-c++": FileType.CODE,
            "text/x-php": FileType.CODE,
            "application/x-php": FileType.CODE,
            "application/javascript": FileType.CODE,
        }
        try:
            mime = magic.Magic(mime=True)
            file_type = mime.from_file(file_path)

            if file_type=="application/zip":
                if self.is_docx(file_path):
                    return FileType.DOCX
            extension = MIME_TO_EXTENSION.get(file_type)
            if not extension:
                logger.error(f"Unsupported file type detected: {file_type}")
                raise Exception("Unsupported file type: This system only supports files.")
            return extension
        except Exception as e:
            logger.error(f"Error detecting file type for {file_path}: {e}")
            raise Exception(f"File type detection failed: {e}")

    def is_docx(self,file_path: str):
        try:
            with zipfile.ZipFile(file_path,'r')as zip_ref:
                return "word/document.xml" in zip_ref.namelist()
        except zipfile.BadZipfile:
                return False

    def append_result(self, results, customer_guid, filename, page_number, content):
        results.append({
            "metadata": {
                "page_number": page_number,
                "customer_guid": customer_guid,
                "filename": filename,
            },
            "text": content
        })

    def format_table_as_text(self, table):
        try:
            if not table:
                logger.error("Empty or None table received for formatting.")
                return "Empty table"

            headers = table[0]  # Assuming the first row is the header
            rows = table[1:]  # All subsequent rows
            result = []

            # Iterate through each row in the table and format the data generically
            for row in rows:
                formatted_row = ", ".join(
                    [f"{headers[i]}: {row[i] if i < len(row) else 'NA'}" for i in range(len(headers))])
                result.append(formatted_row)

            return "\n".join(result).strip()

        except Exception as e:
            logger.error(f"Error formatting table: {e}")
            raise Exception(f"Table formatting failed: {e}")

    def is_within_bbox(self, point: tuple, bbox: tuple):
        x, y = point
        x0, y0, x1, y1 = bbox
        return x0 <= x <= x1 and y0 <= y <= y1

    def is_vertical_text(self, element):
        # Check if both 'x0' and 'x1' keys exist in the element
        if 'x0' in element and 'x1' in element:
            return abs(element["x0"] - element["x1"]) < abs(element["y0"] - element["y1"]) * 0.5
        return False

    def reverse_vertical_text(self, text):
        return text[::-1]

    def extract_pdf_content(self, customer_guid: str, file_path: str, filename: str):

        try:
            results = []
            with pdfplumber.open(file_path) as pdf:
                for page_number, page in enumerate(pdf.pages, start=1):
                    page_elements = []
                    page_width = page.width
                    mid_x = page_width / 2  # Midpoint for left/right classification

                    # Extract text blocks and classify into left/right or full-paragraph
                    left_content = []
                    right_content = []
                    all_content = []

                    for block in page.extract_words():
                        x0 = block["x0"]
                        y0 = block["top"]

                        # Classify into left or right sections
                        if x0 < mid_x:
                            left_content.append({"text": block["text"], "y0": y0})
                        else:
                            right_content.append({"text": block["text"], "y0": y0})

                        all_content.append({"text": block["text"], "y0": y0})

                    # Check if left and right content exists
                    has_left = len(left_content) > 0
                    has_right = len(right_content) > 0

                    # Determine layout type
                    if has_left and has_right:
                        # Left/Right Layout: Process separately and merge
                        left_content = sorted(left_content, key=lambda b: b["y0"])
                        right_content = sorted(right_content, key=lambda b: b["y0"])
                        merged_content = " ".join([b["text"] for b in left_content + right_content])
                    else:
                        # Full-Paragraph Layout: Sort all text blocks and merge them
                        sorted_content = sorted(all_content, key=lambda b: b["y0"])
                        merged_content = " ".join([b["text"] for b in sorted_content])

                    # Append merged text content to page_elements
                    page_elements.append({
                        "type": "text",
                        "content": merged_content,
                        "x0": 0,
                        "y0": 0  # No specific position for merged content
                    })

                    # Extract text and table data from the page
                    tables = page.extract_tables()
                    if tables:
                        for table in tables:
                            # Format and collect table data
                            table_text = self.format_table_as_text(table)
                            page_elements.append({
                                "type": "table",
                                "content": table_text,
                                "x0": 0,
                                "y0": 0
                            })

                    # If there's no extracted text, perform OCR and append image-based text
                    if not page.extract_text():
                        page_image = page.to_image(resolution=300).original
                        ocr_text = pytesseract.image_to_string(page_image)
                        page_elements.append({
                            "type": "image",
                            "content": ocr_text.strip(),
                            "x0": 0,
                            "y0": 0
                        })

                    # Sort page elements by vertical position (y0) and horizontal position (x0)
                    page_elements.sort(key=lambda w: (w["y0"], w["x0"]))

                    # Handle reversed vertical text
                    corrected_content = []
                    for element in page_elements:
                        if element["type"] == "text" and self.is_vertical_text(element):
                            corrected_content.append(self.reverse_vertical_text(element["content"]))
                        else:
                            corrected_content.append(element["content"])

                    # Join the corrected content
                    page_text = " ".join(corrected_content)

                    self.append_result(results,customer_guid,filename,page_number,page_text)

                # Save the extracted content to a file
                output_file_path = f"/tmp/{customer_guid}/{filename}.rawcontent"
                os.makedirs(os.path.dirname(output_file_path), exist_ok=True)
                with open(output_file_path, "w", encoding="utf-8") as raw_content_file:
                    json.dump(results, raw_content_file, indent=4)

                logger.info(f"Extracted content saved to '{output_file_path}'")
                return output_file_path
        except Exception as e:
            logger.error(f"Error extracting content from PDF file '{filename}': {e}")
            raise Exception(f"PDF content extraction failed: {e}")

    def extract_docx_content(self, customer_guid: str, file_path: str, filename: str):
        try:
            results = []
            doc = Document(file_path)

            # Variables for text extraction
            page_elements = []
            left_content = []
            right_content = []
            all_content = []
            simulated_mid_x = 50

            # Track page number and paragraph count
            current_page_number = 1
            current_paragraph_count = 0

            # Process paragraphs for text extraction
            paragraphs = doc.paragraphs
            total_paragraphs = len(paragraphs)
            logger.info(f"Total paragraphs in document: {total_paragraphs}")

            if total_paragraphs == 0:
                logger.error("No paragraphs found in the document!")

            for i, paragraph in enumerate(paragraphs):
                text = paragraph.text.strip()
                if text:
                    logger.info(f"Processing paragraph {i + 1}: {text}")


                    x0 = len(text)
                    y0 = i

                    if x0 < simulated_mid_x:
                        left_content.append({"text": text, "y0": y0})
                    else:
                        right_content.append({"text": text, "y0": y0})

                    all_content.append({"text": text, "y0": y0})
                    current_paragraph_count += 1  # Increment paragraph count

                    if len(all_content) >= 50:
                        logger.info(
                            f"Switching to page {current_page_number + 1} after processing {current_paragraph_count} paragraphs.")

                        # Sort and merge left/right content by y0
                        left_content = sorted(left_content, key=lambda b: b["y0"])
                        right_content = sorted(right_content, key=lambda b: b["y0"])
                        merged_content = " ".join([b["text"] for b in left_content + right_content])

                        # Add the content for the current page to the result
                        page_elements.append({
                            "metadata": {"page_number": current_page_number},
                            "type": "text",
                            "content": merged_content,
                            "x0": 0,
                            "y0": 0
                        })

                        # Add the page's content to results
                        page_text = " ".join(
                            [element["content"] for element in page_elements if element["type"] == "text"])
                        results.append({
                            "metadata": {"page_number": current_page_number},
                            "text": page_text
                        })

                        # Reset for the next page
                        left_content = []
                        right_content = []
                        all_content = []
                        page_elements = []  # Reset page elements after each page
                        current_paragraph_count = 0  # Reset paragraph count for the new page
                        current_page_number += 1  # Increment page number

            # After the loop, handle the last page (if any content remains)
            if current_paragraph_count > 0:
                # Sort and merge left/right content
                left_content = sorted(left_content, key=lambda b: b["y0"])
                right_content = sorted(right_content, key=lambda b: b["y0"])
                merged_content = " ".join([b["text"] for b in left_content + right_content])

                # Add the last page's content to the result
                page_elements.append({
                    "metadata": {"page_number": current_page_number},
                    "type": "text",
                    "content": merged_content,
                    "x0": 0,
                    "y0": 0
                })

            # Extract tables and format them as text
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)

                table_text = self.format_table_as_text(table_data)
                page_elements.append({
                    "type": "table",
                    "content": table_text,
                    "x0": 0,
                    "y0": 0
                })

            # Handle embedded images and extract text using OCR
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_data = rel.target_part.blob
                    image = Image.open(io.BytesIO(image_data))
                    ocr_text = pytesseract.image_to_string(image)

                    page_elements.append({
                        "type": "image",
                        "content": ocr_text.strip(),
                        "x0": 0,
                        "y0": 0
                    })

            # Sort all page elements (text, tables, and images) by vertical (y0) and horizontal (x0) position
            page_elements.sort(key=lambda w: (w["y0"], w["x0"]))

            # Prepare the corrected content (combine text, tables, and OCR from images)
            corrected_content = []
            for element in page_elements:
                if element["type"] == "table":
                    corrected_content.append(element["content"])  # Table text is pre-formatted
                elif element["type"] == "image":
                    corrected_content.append(element["content"])  # OCR text
                else:
                    corrected_content.append(element["content"])  # Text content

            # Combine the corrected content into a single text string
            page_text = " ".join(corrected_content)

            # Append the final content to results
            self.append_result(results,customer_guid,filename,current_page_number,page_text)

            # Save the extracted content to a file
            output_file_path = f"/tmp/{customer_guid}/{filename}.rawcontent"
            os.makedirs(os.path.dirname(output_file_path), exist_ok=True)
            with open(output_file_path, "w", encoding="utf-8") as raw_content_file:
                json.dump(results, raw_content_file, indent=4)

            logger.info(f"Extracted content saved to '{output_file_path}'")
            return output_file_path

        except Exception as e:
            logger.error(f"Error extracting content from DOCX file '{filename}': {e}")
            raise Exception(f"DOCX content extraction failed: {e}")

    def extract_ppt_content(self, customer_guid: str, file_path: str, filename: str):
        try:

            results = []
            presentation = Presentation(file_path)

            for slide_number, slide in enumerate(presentation.slides, start=1):
                slide_elements = []

                # Extract slide text
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_elements.append({
                                    "type": "text",
                                    "content": text
                                })

                    # Extract table data
                    if shape.has_table:
                        table = shape.table
                        table_data = [
                            [cell.text.strip() for cell in row.cells]
                            for row in table.rows
                        ]
                        table_text = self.format_table_as_text(table_data)
                        slide_elements.append({
                            "type": "table",
                            "content": table_text
                        })

                    # Extract images
                    if shape.shape_type == CustomShapeType.PICTURE.value:  # 13 indicates a picture
                        image = shape.image
                        image_data = image.blob
                        image_obj = Image.open(io.BytesIO(image_data))
                        ocr_text = pytesseract.image_to_string(image_obj).strip()
                        slide_elements.append({
                            "type": "image",
                            "content": ocr_text
                        })
                        
                    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        chart = shape.chart
                        table_data = []
                        x_axis_labels = [category for category in chart.plots[0].categories]
                        first_column_name = chart.series[0].name if chart.series and chart.series[0].name else None
                        series_names = []

                        for series in chart.series:
                            series_name = series.name if series.name else "Unnamed Series"
                            series_names.append(series_name)
                        table_data.append([first_column_name] + series_names)
                        for idx, category in enumerate(x_axis_labels):
                            row = [category]
                            for series in chart.series:
                                value = series.values[idx] if idx < len(series.values) else None
                                row.append(value)
                            table_data.append(row)

                        # Append the reconstructed table data
                        slide_elements.append({
                            "type": "chart_table",
                            "content": table_data
                        })   

                # Combine slide elements
                slide_text = "\n".join(
                    json.dumps(element["content"]) if isinstance(element["content"], list) else element["content"]
                    for element in slide_elements
                )
                self.append_result(results,customer_guid,filename,slide_number,slide_text)

            # Save results to a file
            output_file_path = f"/tmp/{customer_guid}/{filename}.rawcontent"
            os.makedirs(os.path.dirname(output_file_path), exist_ok=True)
            with open(output_file_path, "w", encoding="utf-8") as raw_content_file:
                json.dump(results, raw_content_file, indent=4)

            logger.info(f"Extracted content saved to '{output_file_path}'")
            return output_file_path

        except Exception as e:
            logger.error(f"Error extracting content from PPT file '{filename}': {e}")
            raise Exception(f"PPTX content extraction failed: {e}")
            
    def extract_excel_content(self, customer_guid: str, file_path: str, filename: str):
        try:
            results = []
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            chart_types = (
                BarChart, PieChart, LineChart, AreaChart, RadarChart, ScatterChart,
                BubbleChart, DoughnutChart, StockChart, SurfaceChart
            )

            for sheet_index, sheet_name in enumerate(workbook.sheetnames, start=1):
                sheet = workbook[sheet_name]
                sheet_elements = []

                table_data = []
                for row in sheet.iter_rows():
                    row_data = [cell.value if cell.value is not None else "" for cell in row]
                    table_data.append(row_data)

                if table_data:
                    formatted_table = self.format_table_as_text(table_data)
                    sheet_elements.append({
                        "type": "table",
                        "content": formatted_table
                    })

                if hasattr(sheet, '_images') and sheet._images:
                    for image in sheet._images:
                        if isinstance(image, OpenpyxlImage):
                            image_stream = image.ref
                            image_obj = Image.open(image_stream)
                            ocr_text = pytesseract.image_to_string(image_obj).strip()
                            sheet_elements.append({
                                "type": "image",
                                "content": ocr_text
                            })

                if hasattr(sheet, '_charts'):
                    for drawing in sheet._charts:
                        if isinstance(drawing, chart_types):
                            chart_info = {
                                'type': drawing.__class__.__name__,
                                'title': self.get_text(drawing.title) if drawing.title else None,
                                'data': []
                            }

                            for index, series in enumerate(drawing.series):
                                series_data = {
                                    'name': self.get_series_name(index),
                                    'values': self.get_values_from_series(series)
                                }
                                chart_info['data'].append(series_data)

                            sheet_elements.append({
                                "type": "chart",
                                "content": chart_info
                            })

                # Combine sheet content
                combined_text = "\n".join(
                    json.dumps(element["content"], indent=2) if isinstance(element["content"], (list, dict)) else
                    element["content"]
                    for element in sheet_elements
                )

                self.append_result(results,customer_guid,filename,sheet_index,combined_text)

            # Save results to a file
            output_file_path = f"/tmp/{customer_guid}/{filename}.rawcontent"
            os.makedirs(os.path.dirname(output_file_path), exist_ok=True)
            with open(output_file_path, "w", encoding="utf-8") as raw_content_file:
                json.dump(results, raw_content_file, indent=4)

            logger.info(f"Extracted content saved to '{output_file_path}'")
            return output_file_path

        except Exception as e:
            logger.error(f"Error extracting content from Excel file '{filename}': {e}")
            raise Exception(f"Excel content extraction failed: {e}")

    def get_series_name(self, index):
        return f"Series {index + 1}"

    def get_values_from_series(self, series):
        try:
            if hasattr(series, 'values') and isinstance(series.values, openpyxl.chart.Reference):
                return [cell.value for cell in series.values.cells]
            else:
                return None
        except Exception as e:
            logger.warning(f"Failed to extract values from series: {e}")
            return None

    def get_text(self, obj):
        if obj is None:
            return None
        elif isinstance(obj, str):
            return obj
        elif isinstance(obj, openpyxl.chart.text.Text):
            return self.extract_text_from_text_object(obj)
        elif isinstance(obj, openpyxl.chart.text.RichText):
            return self.extract_text_from_richtext_object(obj)
        else:
            return str(obj)

    def extract_html_content(self, customer_guid: str, source: str, filename: str):
        try:
            results = []

            # Read the HTML file from the local file system
            logger.debug(f"Reading local HTML file from: {source}")
            with open(source, "r", encoding="utf-8") as html_file:
                html_content = html_file.read()

            # Parse the HTML content
            soup = BeautifulSoup(html_content, "html.parser")

            # Identify pages or fallback to treat the entire HTML as one page
            pages = soup.find_all("div", class_="page")
            if not pages:
                pages = [soup]

            page_number = 0
            for page in pages:
                page_number += 1
                logger.debug(f"Processing page {page_number}")
                text_content = []

                for text_element in page.find_all(string=True):
                    if isinstance(text_element, Comment) or text_element.parent.name in ["script", "style", "meta",
                                                                                         "[document]"]:
                        continue

                    clean_text = text_element.strip()
                    text_content.append(clean_text)

                # Extract tables from the current page
                tables = page.find_all("table")
                if tables:
                    logger.debug(f"Found {len(tables)} tables on page {page_number}")
                    for table in tables:
                        headers = []
                        rows = []

                        # Extract table headers
                        header_row = table.find("tr")
                        if header_row:
                            headers = [header.get_text(strip=True) for header in header_row.find_all("th")]

                        # Extract table rows
                        for row in table.find_all("tr")[1:]:  # Skip the header row
                            row_data = [cell.get_text(strip=True) for cell in row.find_all("td")]
                            if row_data:
                                rows.append(row_data)

                        # Format the table as text
                        if headers or rows:
                            formatted_table = self.format_table_as_text([headers] + rows)
                            text_content.append(formatted_table)

                # Extract charts and their data from script tags
                scripts = page.find_all("script")
                for script in scripts:
                    script_content = script.string
                    if script_content and "new Chart" in script_content:
                        # Regex to find chart data from script
                        chart_labels = re.findall(r"labels:\s*\[(.*?)\]", script_content)
                        chart_values = re.findall(r"data:\s*\[(.*?)\]", script_content)

                        for chart_label, chart_value in zip(chart_labels, chart_values):
                            labels = [label.strip().replace("'", "") for label in chart_label.split(",")]
                            data = [int(x.strip()) for x in chart_value.split(",")]

                            # Format the chart data as plain text
                            formatted_chart = self.format_chart_as_table(labels, data)
                            text_content.append(formatted_chart)
                            logger.debug(f"Chart found with labels: {labels}")
                            logger.debug(f"Chart data: {data}")

                # Extract images and perform OCR
                images = page.find_all("img")
                for img in images:
                    img_src = img.get("src")
                    if img_src:
                        try:
                            # Resolve the image path (local or remote)
                            image_path = os.path.join(os.path.dirname(source), img_src) if not img_src.startswith(
                                "http") else img_src
                            image_data = None

                            if img_src.startswith("http"):
                                response = requests.get(image_path)
                                response.raise_for_status()
                                image_data = Image.open(BytesIO(response.content))
                            else:
                                if os.path.exists(image_path):
                                    image_data = Image.open(image_path)
                                else:
                                    logger.error(f"Image file not found: {image_path}")
                                    continue

                            # OCR to extract text from image
                            ocr_text = pytesseract.image_to_string(image_data).strip()
                            if ocr_text:
                                text_content.append(ocr_text)
                        except Exception as e:
                            logger.error(f"Failed to process image {img_src}: {e}")

                            # Append footer content to the text_content
                            footer = soup.find("footer")  # Look for a <footer> tag
                            if footer:
                                for text_element in footer.find_all(string=True):
                                    # Skip comments and other non-visible elements
                                    if isinstance(text_element, Comment) or text_element.parent.name in ["script",
                                                                                                         "style",
                                                                                                         "meta",
                                                                                                         "[document]"]:
                                        continue

                                    clean_text = text_element.strip()
                                    if clean_text:
                                        text_content.append(clean_text)
                content = " ".join(text_content)
                # Append the result for the current page
                self.append_result(results,customer_guid,filename,page_number,content)

            # Save extracted content
            output_file_path = f"/tmp/{customer_guid}/{filename}.rawcontent"
            os.makedirs(os.path.dirname(output_file_path), exist_ok=True)
            with open(output_file_path, "w", encoding="utf-8") as raw_content_file:
                json.dump(results, raw_content_file, indent=4)

            logger.info(f"Extracted HTML content saved to '{output_file_path}'")
            return output_file_path

        except Exception as e:
            logger.error(f"Error extracting content from HTML file '{filename}': {e}")
            raise Exception(f"HTML content extraction failed: {e}")

    def format_chart_as_table(self, labels, data):
        try:
            if not labels or not data:
                logger.error("Empty or invalid chart data received for formatting.")
                return "Empty chart data"

            # Format chart data into a table-like structure
            headers = ["Label", "Data"]
            rows = zip(labels, data)

            result = []
            for row in rows:
                formatted_row = ", ".join(
                    [f"{headers[i]}: {row[i]}" for i in range(len(headers))]
                )
                result.append(formatted_row)

            return "\n".join(result).strip()

        except Exception as e:
            logger.error(f"Error formatting chart data: {e}")
            raise Exception(f"Chart formatting failed: {e}")

    def extract_json_content(self, customer_guid: str, file_path: str, filename: str):
        try:
            results = []
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
                extracted_data = self.extract_dynamic_json(data)  # Updated function call
                structured_data = self.extract_data_from_json(extracted_data)
                formatted_data = self.format_output(structured_data)
                if not formatted_data.strip():
                    return None
                self.append_result(results,customer_guid,filename,page_number=1,content=formatted_data)
            # Save the extracted content to a file
            output_file_path = f"/tmp/{customer_guid}/{filename}.rawcontent"
            os.makedirs(os.path.dirname(output_file_path), exist_ok=True)
            with open(output_file_path, "w", encoding="utf-8") as raw_content_file:
                json.dump(results, raw_content_file, indent=4)
            return output_file_path
        except Exception as e:
            raise Exception(f"JSON content extraction failed: {e}")

    def extract_dynamic_json(self, data, path=""):
        extracted_content = []
        if isinstance(data, dict):
            for key, value in data.items():
                full_path = f"{path}.{key}" if path else key
                if isinstance(value, (dict, list)):
                    extracted_content.extend(self.extract_dynamic_json(value, full_path))
                else:
                    extracted_content.append(f"{full_path}: {value}")
        elif isinstance(data, list):
            if not data:
                print(f"Warning: Empty list at {path}")
            for idx, item in enumerate(data):
                full_path = f"{path}[{idx}]"
                if isinstance(item, (dict, list)):
                    extracted_content.extend(self.extract_dynamic_json(item, full_path))
                else:
                    extracted_content.append(f"{full_path}: {item}")
        return extracted_content

    def extract_data_from_json(self, json_data):
        structured_data = {}
        if isinstance(json_data, list):
            for item in json_data:
                if isinstance(item, str):
                    if ": " in item:
                        path, value = item.split(": ", 1)
                        path_parts = path.split('.')
                        temp = structured_data
                        for part in path_parts[:-1]:
                            temp = temp.setdefault(part, {})
                        temp[path_parts[-1]] = value
        return structured_data

    def format_output(self, extracted_data):
        formatted_output = []

        def format_section(data, section_name=""):
            nonlocal formatted_output
            if isinstance(data, dict):
                if section_name:
                    formatted_output.append(f"{section_name.capitalize()}")
                for key, value in data.items():
                    if isinstance(value, (dict, list)):
                        format_section(value, key)
                    else:
                        formatted_output.append(f"    {key}: {value}")
            elif isinstance(data, list):
                for idx, item in enumerate(data):
                    format_section(item, f"{section_name}[{idx}]")
            else:
                formatted_output.append(f"    {section_name}: {data}")

        format_section(extracted_data)
        return "\n".join(formatted_output)

    def extract_image_content(self, customer_guid: str, image_path: str, filename: str):

        try:
            logger.info(f"Extracting text from image '{image_path}' for customer '{customer_guid}'")
            image = Image.open(image_path)

            ocr_text = pytesseract.image_to_string(image)

            results = []

            corrected_code = ocr_text.strip()

            self.append_result(results,customer_guid,filename,page_number=1,content=corrected_code)

            output_file_path = f"/tmp/{customer_guid}/{filename}.rawcontent"
            os.makedirs(os.path.dirname(output_file_path), exist_ok=True)

            # Save the extracted text to a JSON file
            with open(output_file_path, "w", encoding="utf-8") as raw_content_file:
                json.dump(results, raw_content_file, indent=4)

            logger.info(f"Extracted content saved to '{output_file_path}'")
            return output_file_path
        except Exception as e:
            logger.error(f"Error extracting text from image '{filename}': {e}")
            raise Exception(f"Image extraction failed: {e}")

    def extract_yaml_content(self, customer_guid: str, file_path: str, filename: str):
        try:
            results = []
            with open(file_path, 'r', encoding='utf-8') as file:
                yaml_data = yaml.safe_load(file)  # Parse YAML content
                formatted_data = self.yaml_format_output(yaml_data)  # Format YAML content
                if not formatted_data.strip():
                    return None
                self.append_result(results,customer_guid,filename,page_number=1,content=formatted_data)

            # Save the extracted content to a file
            output_file_path = f"/tmp/{customer_guid}/{filename}.rawcontent"
            os.makedirs(os.path.dirname(output_file_path), exist_ok=True)
            with open(output_file_path, "w", encoding="utf-8") as raw_content_file:
                json.dump(results, raw_content_file, indent=4)
            return output_file_path

        except Exception as e:
            raise Exception(f"YAML content extraction failed: {e}")

    def yaml_format_output(self, data, indent=0):
        formatted_output = []

        def format_yaml(data, current_indent):
            if isinstance(data, dict):
                for key, value in data.items():
                    if isinstance(value, (dict, list)):
                        formatted_output.append(f"{' ' * current_indent}{key}:")
                        format_yaml(value, current_indent + 2)
                    else:
                        # Preserve multiline strings with proper YAML syntax
                        if isinstance(value, str) and '\n' in value:
                            formatted_output.append(f"{' ' * current_indent}{key}: |")
                            for line in value.splitlines():
                                formatted_output.append(f"{' ' * (current_indent + 2)}{line}")
                        else:
                            formatted_output.append(f"{' ' * current_indent}{key}: {value}")
            elif isinstance(data, list):
                for item in data:
                    if isinstance(item, (dict, list)):
                        formatted_output.append(f"{' ' * current_indent}-")
                        format_yaml(item, current_indent + 2)
                    else:
                        formatted_output.append(f"{' ' * current_indent}- {item}")
            else:
                formatted_output.append(f"{' ' * current_indent}{data}")

        format_yaml(data, indent)
        return "\n".join(formatted_output)

    def extract_code_file_content(self, customer_guid: str, file_path: str, filename: str):
        try:
            results = []
            # Read the file content
            with open(file_path, "r", encoding="utf-8") as file:
                code_content = file.read()

            self.append_result(results,customer_guid,filename,page_number=1,content=code_content)

            # Save the extracted content to a file
            output_file_path = f"/tmp/{customer_guid}/{filename}.rawcontent"
            os.makedirs(os.path.dirname(output_file_path), exist_ok=True)
            with open(output_file_path, "w", encoding="utf-8") as raw_content_file:
                json.dump(results, raw_content_file, indent=4)
            logger.info(f"Extracted code content saved to '{output_file_path}'")
            return output_file_path

        except Exception as e:
            logger.error(f"Code content extraction failed for file '{filename}': {e}")
            raise Exception(f"Code content extraction failed: {e}")


if __name__ == "__main__":
    customer_guid = "d3d62569-03b4-41ce-94bf-918a97684763"
    url="https://www.bookwidgets.com/blog/2016/10/15-education-blogs-every-teacher-should-know-about"
    filename = "bookwidgets.html"
    upload_file_for_chunks = UploadFileForChunks()
    upload_file_for_chunks.extract_html_files(customer_guid,url, filename)
